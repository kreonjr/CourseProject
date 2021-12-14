# -*- coding: utf-8 -*-
"""
get_project_text.py

Pulls contents of textual files from locally-cloned GitHub projects.
Handles .md, .pdf, .docx, .pptx.
Stores text in two output files: one with a single row per project file,
and another with a single row per project.

Arguments:
    --projectlist: CSV list of locally-cloned GitHub projects.
        Defaults to "repo_forks.csv" in the same folder as this script.
    --projectroot: Root directory of locally-cloned GitHub projects.
        Defaults to subfolder "repo_forks" of the directory this script is in.
    --outputdir: Destination directory of output files.
        Defaults to the folder this script is in.

Created on Sat Oct 30 20:55:24 2021

@author: ktuoh
"""

import argparse
import os
import re
from unidecode import unidecode
import pandas as pd
import fitz
import docx2txt
from pptx import Presentation
from bs4 import BeautifulSoup
from markdown import markdown
import time
import threading

# Arguments:
    # projectlist: CSV list of projects.
    # outputdir: destination directory of output files.

# Steps:
    # Read in list of projects
    # For each project:
        # Look for textual files
        # For each file:
            # Get name
            # Get contents
            # Normalize whitespace
            # Concatenate contents into single line
            # Append project name, file name, and contents to data frame
    # Output project/file contents
    # Aggregate project/file contents into project contents
    # Output project contents

def get_pdf_text(filepath):
    """
    Return text contents of the given PDF.
    Return as string.
    """

    filetext = ""
    
    try:
        pdf_doc = fitz.open(filepath)

    # Handle incorrectly-formatted PDFs
    except Exception as e:
        print("Error opening PDF file", filepath, e)

    # Pull text from readable PDFs
    else:
        # Loop through PDF pages and collect text.
    
        for pdf_page in pdf_doc:
            filetext = filetext + pdf_page.get_text()
            
        pdf_doc.close()
    
    return filetext


def get_docx_text(filepath):
    """
    Return text of the given Word document.
    Return as string.
    """
    
    filetext = ""

    try:
        filetext = docx2txt.process(filepath)
    except Exception as e:
        print("Error processing Word document", filepath, e)
    
    return filetext
    

def get_pptx_text(filepath):
    """
    Return text of the given Powerpoint document.
    Return as string.
    """
    
    filetext = ""
    
    try:
        ppt = Presentation(filepath)
    
        # Loop through shapes in slides and collect any text.
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    filetext = filetext + " " + shape.text

    except Exception as e:
        print("Error processing Powerpoint file", filepath, e)
                

    return filetext

def get_md_text(filepath):
    """
    Return flattened string consisting of the contents
    of the given markdown file (or any text file).
    """
    
    filetext = ""
    
    try:
        fd = open(filepath, "r", encoding="utf-8", errors="ignore")
        md_text = fd.read()
        fd.close()
        
        md_html = markdown(md_text)
        filetext = "".join(BeautifulSoup(md_html, "html.parser").findAll(text=True))
        
    except Exception as e:
        print("Error processing markdown file", filepath, e)
    
    return filetext
        
    

def normalize_text(input_text):
    """
    Given a string, remove or replace non-ASCII characters,
    normalize spaces, and make text a single line.
    """
    
    # Reduce text to single line.
    norm_text = " ".join(input_text.splitlines()).strip()

    # Normalize whitespace
    norm_text = re.sub("\s+", " ", norm_text)

    # Convert Unicode characters to their ASCII equivalents and return.
    return unidecode(norm_text)


def process_project_files(project_url, project_folder):
    """
    Given a project folder, identify textual files and pull contents.
    Return the name of each file, plus the text as a single line.
    Look in top level of directory only.
    """
    
    filelist = []

    for filename in os.listdir(project_folder):
        filepath = os.path.join(project_folder, filename)
        if os.path.isfile(filepath):
            filetext = ""
            if (filename.endswith(".pdf")):
                filetext = get_pdf_text(filepath)
            elif (filename.endswith(".docx") and not filename.startswith("~$")):
                filetext = get_docx_text(filepath)
            elif (filename.endswith(".pptx") and not filename.startswith("~$")):
                filetext = get_pptx_text(filepath)
            elif (filename.endswith(".md")):
                filetext = get_md_text(filepath)
                
            if (filetext != ""):
                # Remove or replace non-ASCII characters, normalize spaces, and make text a single line
                normtext = normalize_text(filetext)
                file_dict = {
                    "project_url": project_url,
                    "file_name": filename,
                    "file_text": normtext
                    }
                filelist.append(file_dict)
    
    return filelist


def process_project_list(project_df, project_root):
    """
    Iterate through dataframe of info about GitHub projects.
    Look for textual files in each project, read the text, and log it.
    project_list: Pandas dataframe
    project_root: root directory of locally-cloned projects.
    """
    
    doclist = []
    
    for index, row in project_df.iterrows():
        project_folder = os.path.join(project_root, row["owner_login"])
        if (os.path.isdir(project_folder)):
            # Call routine to look for textual files in project_folder and collect their contents
            doclist = doclist + process_project_files(row["project_url"], project_folder)
            
    return pd.DataFrame(doclist)


def main():
    done = False
    def animate():
        bar = [
            " [=     ]",
            " [ =    ]",
            " [  =   ]",
            " [   =  ]",
            " [    = ]",
            " [     =]",
            " [    = ]",
            " [   =  ]",
            " [  =   ]",
            " [ =    ]",
        ]
        i = 0

        while not done:
            print("Extracting text data", bar[i % len(bar)], end="\r")
            time.sleep(.2)
            i += 1

    t = threading.Thread(target=animate)
    t.start()

    # Get folder this script is running from
    script_dir = os.path.dirname(os.path.abspath(__file__)).replace("\\", "/")
    
    # Get arguments
    parser = argparse.ArgumentParser()
    parser.add_argument("--projectlist", help="CSV list of locally-cloned GitHub projects", default=os.path.join(script_dir, "repo_forks.csv"))
    parser.add_argument("--projectroot", help="Root directory of locally-cloned GitHub projects", default=os.path.join(script_dir, "repo_forks"))
    parser.add_argument("--outputdir", help="Destination directory of output", default=script_dir)
    args = parser.parse_args()

    # Read project list
    try:
        project_df = pd.read_csv(args.projectlist)
    except Exception:
        print("Error reading project list")
        exit()
        
    # Get dataframe of projects, textual files and their contents. One row per file.
    project_file_df = process_project_list(project_df, args.projectroot)
    # Output results
    project_file_df.to_csv(os.path.join(args.outputdir, "project_file_text.tsv"), sep="\t", index=False)

    # Concatenate text of each project's files and produce one row per project
    project_text_series = project_file_df.groupby("project_url")["file_text"].agg(" ".join)
    project_text_series.to_csv(os.path.join(args.outputdir, "project_text.tsv"), sep="\t", index=True)

    done = True
    

if __name__ == "__main__":
    main()
    